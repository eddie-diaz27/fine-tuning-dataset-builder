"""
File processor for extracting text from various file formats.
Supports: PDF, DOCX, TXT, CSV, XLSX, PPTX, MD, HTML
"""

import os
from pathlib import Path
from typing import Optional
from dataclasses import dataclass
from datetime import datetime
import chardet


class FileProcessingError(Exception):
    """Raised when file processing fails"""
    pass


@dataclass
class ProcessedFile:
    """Processed file data"""
    text: str
    file_name: str
    file_type: str
    file_size: int
    char_count: int
    created_at: datetime


class FileProcessor:
    """Process various file formats into text"""

    SUPPORTED_FORMATS = {
        '.pdf': 'extract_pdf',
        '.docx': 'extract_docx',
        '.txt': 'extract_text',
        '.csv': 'extract_csv',
        '.xlsx': 'extract_xlsx',
        '.pptx': 'extract_pptx',
        '.md': 'extract_markdown',
        '.html': 'extract_html'
    }

    def __init__(self, max_size_mb: int = 50, encoding: str = 'utf-8'):
        """
        Initialize file processor.

        Args:
            max_size_mb: Maximum file size in MB
            encoding: Default text encoding
        """
        self.max_size_mb = max_size_mb
        self.encoding = encoding

    def process_file(self, file_path: str) -> ProcessedFile:
        """
        Extract text from file.

        Args:
            file_path: Path to file

        Returns:
            ProcessedFile object with extracted text

        Raises:
            FileProcessingError: If file cannot be processed
        """
        path = Path(file_path)

        # Check if file exists
        if not path.exists():
            raise FileProcessingError(f"File not found: {file_path}")

        # Check file size
        file_size = path.stat().st_size
        file_size_mb = file_size / (1024 * 1024)

        if file_size_mb > self.max_size_mb:
            raise FileProcessingError(
                f"File too large: {file_size_mb:.1f}MB "
                f"(limit: {self.max_size_mb}MB)"
            )

        # Check if empty
        if file_size == 0:
            raise FileProcessingError(f"File is empty: {file_path}")

        # Get file extension
        ext = path.suffix.lower()

        if ext not in self.SUPPORTED_FORMATS:
            raise FileProcessingError(
                f"Unsupported file format: {ext}\n"
                f"Supported formats: {', '.join(self.SUPPORTED_FORMATS.keys())}"
            )

        # Extract text using appropriate method
        method_name = self.SUPPORTED_FORMATS[ext]
        extract_method = getattr(self, method_name)

        try:
            text = extract_method(str(path))
        except Exception as e:
            raise FileProcessingError(
                f"Failed to process {path.name}: {str(e)}"
            )

        # Create ProcessedFile object
        return ProcessedFile(
            text=text,
            file_name=path.name,
            file_type=ext[1:],  # Remove leading dot
            file_size=file_size,
            char_count=len(text),
            created_at=datetime.now()
        )

    def extract_pdf(self, path: str) -> str:
        """
        Extract text from PDF file.

        Args:
            path: Path to PDF file

        Returns:
            Extracted text
        """
        try:
            import PyPDF2
        except ImportError:
            raise FileProcessingError(
                "PyPDF2 is required for PDF processing.\n"
                "Install with: pip install PyPDF2"
            )

        try:
            text_chunks = []

            with open(path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)

                # Check if encrypted
                if reader.is_encrypted:
                    raise FileProcessingError(
                        "PDF is password-protected. "
                        "Please decrypt the file first."
                    )

                # Extract text page by page (memory efficient)
                for page_num, page in enumerate(reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text.strip():
                            text_chunks.append(page_text)
                    except Exception as e:
                        # Continue even if one page fails
                        print(f"Warning: Failed to extract page {page_num + 1}: {e}")

            return "\n\n".join(text_chunks)

        except PyPDF2.errors.PdfReadError as e:
            raise FileProcessingError(f"Corrupted PDF file: {e}")

    def extract_docx(self, path: str) -> str:
        """
        Extract text from DOCX file.

        Args:
            path: Path to DOCX file

        Returns:
            Extracted text
        """
        try:
            from docx import Document
        except ImportError:
            raise FileProcessingError(
                "python-docx is required for DOCX processing.\n"
                "Install with: pip install python-docx"
            )

        try:
            doc = Document(path)
            text_chunks = []

            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_chunks.append(para.text)

            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(cell.text for cell in row.cells)
                    if row_text.strip():
                        text_chunks.append(row_text)

            return "\n\n".join(text_chunks)

        except Exception as e:
            raise FileProcessingError(f"Failed to read DOCX file: {e}")

    def extract_text(self, path: str) -> str:
        """
        Extract text from plain text file.

        Args:
            path: Path to text file

        Returns:
            File contents
        """
        # Try specified encoding first
        try:
            with open(path, 'r', encoding=self.encoding) as f:
                return f.read()
        except UnicodeDecodeError:
            # Detect encoding
            with open(path, 'rb') as f:
                raw_data = f.read()
                result = chardet.detect(raw_data)
                detected_encoding = result['encoding']

            if detected_encoding:
                try:
                    with open(path, 'r', encoding=detected_encoding) as f:
                        return f.read()
                except Exception:
                    pass

            # Fallback to latin-1 (never fails)
            with open(path, 'r', encoding='latin-1') as f:
                return f.read()

    def extract_csv(self, path: str) -> str:
        """
        Extract text from CSV file.

        Args:
            path: Path to CSV file

        Returns:
            Extracted text
        """
        import csv

        try:
            text_chunks = []

            with open(path, 'r', encoding=self.encoding, newline='') as f:
                reader = csv.reader(f)

                # Read all rows
                for row in reader:
                    row_text = ' | '.join(str(cell) for cell in row if cell)
                    if row_text.strip():
                        text_chunks.append(row_text)

            return "\n".join(text_chunks)

        except Exception as e:
            raise FileProcessingError(f"Failed to read CSV file: {e}")

    def extract_xlsx(self, path: str) -> str:
        """
        Extract text from XLSX file.

        Args:
            path: Path to XLSX file

        Returns:
            Extracted text
        """
        try:
            import openpyxl
        except ImportError:
            raise FileProcessingError(
                "openpyxl is required for XLSX processing.\n"
                "Install with: pip install openpyxl"
            )

        try:
            workbook = openpyxl.load_workbook(path, data_only=True)
            text_chunks = []

            # Process each sheet
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]

                # Add sheet name as header
                text_chunks.append(f"=== {sheet_name} ===")

                # Extract rows
                for row in sheet.iter_rows(values_only=True):
                    row_text = ' | '.join(str(cell) for cell in row if cell is not None)
                    if row_text.strip():
                        text_chunks.append(row_text)

                text_chunks.append("")  # Blank line between sheets

            return "\n".join(text_chunks)

        except Exception as e:
            raise FileProcessingError(f"Failed to read XLSX file: {e}")

    def extract_pptx(self, path: str) -> str:
        """
        Extract text from PPTX file.

        Args:
            path: Path to PPTX file

        Returns:
            Extracted text
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise FileProcessingError(
                "python-pptx is required for PPTX processing.\n"
                "Install with: pip install python-pptx"
            )

        try:
            prs = Presentation(path)
            text_chunks = []

            # Process each slide
            for slide_num, slide in enumerate(prs.slides, 1):
                text_chunks.append(f"=== Slide {slide_num} ===")

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_chunks.append(shape.text)

                text_chunks.append("")  # Blank line between slides

            return "\n".join(text_chunks)

        except Exception as e:
            raise FileProcessingError(f"Failed to read PPTX file: {e}")

    def extract_markdown(self, path: str) -> str:
        """
        Extract text from Markdown file.

        Args:
            path: Path to MD file

        Returns:
            File contents (markdown preserved)
        """
        return self.extract_text(path)

    def extract_html(self, path: str) -> str:
        """
        Extract text from HTML file.

        Args:
            path: Path to HTML file

        Returns:
            Extracted text
        """
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise FileProcessingError(
                "beautifulsoup4 is required for HTML processing.\n"
                "Install with: pip install beautifulsoup4 lxml"
            )

        try:
            with open(path, 'r', encoding=self.encoding) as f:
                html_content = f.read()

            soup = BeautifulSoup(html_content, 'lxml')

            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()

            # Get text
            text = soup.get_text()

            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)

            return text

        except Exception as e:
            raise FileProcessingError(f"Failed to read HTML file: {e}")
